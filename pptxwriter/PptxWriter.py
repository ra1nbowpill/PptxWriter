import os
import logging
from math import ceil
from itertools import product
from functools import partial

import numpy as np
import pandas as pd
from unidecode import unidecode

# Powerpoint Generation
# import pptx
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.dml import MSO_THEME_COLOR as COLOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE as CHART
from pptx.enum.chart import XL_LEGEND_POSITION as LEGEND_POS
"""from pptx.enum.chart import XL_LABEL_POSITION as DL_POS
from pptx.enum.chart import XL_LABEL_POSITION as LABEL_POS"""
from pptx.enum.chart import XL_TICK_LABEL_POSITION as TICK_LABEL_POS

# Write a PowerPoint Table using a Panda's DataFrame
import pd2pptx
from pd2pptx import TableStyle

logging.basicConfig(level=logging.INFO)

"""

    TODO:
    - make a package for PptxWriter class
    - to the TODO in add_chart_slide
    - cursor

"""


def capitalize(string):
    """Capitalize a string

    Make the first character of a string uppercase
    This differs from str.capitalize as it lower the other character
    :param string: string to capitalize
    :type string: str
    :returns: capitalized string
    :rtype: {str}
    """
    return string[0].upper() + string[1:]


def to_ascii(string):
    """Convert a string to ASCII and remove non ASCII characters

    [description]
    :param string: string to be clean
    :type string: str
    :returns: string without non ASCII characters
    :rtype: {str}
    """
    return string.encode('ascii', 'ignore').decode('utf-8')


def clean_title(string):
    """Clean a string to be used as filename

    :param string: string to clean
    :type string: str
    :returns: string to be used as a filename
    :rtype: {str}
    """
    string = unidecode(string)
    string = string.strip()
    string = string.replace('/', '')
    string = string.replace('?', '')
    # Common character in excel/powerpoint
    string = string.replace('\\xa0', ' ')
    string = string.replace(' ', '_')
    string = to_ascii(string)
    return string


def uniform_slice(iterable, max_slice_size):
    """Split an iterable in uniform sized parts

    Yield iterable's part of uniform size (the max number of element
    in one part will be `max_slice_part`)
    :param iterable: Iterable to split
    :type iterable: iterable
    :param max_batch: [description]
    :type max_batch: int > 0
    :returns: Batch of the iterable
    :rtype: {list iterable}
    """
    if max_slice_size <= 0:
        raise ValueError('max_slice_size must be > to 0')

    len_iter = len(iterable)
    to_treat = len_iter
    treated = 0
    nb_slice = ceil(len_iter / max_slice_size)

    while nb_slice != 0:
        nb_treated = ceil(to_treat / nb_slice)
        nb_slice -= 1
        yield iterable[treated:treated + nb_treated]
        treated += nb_treated
        to_treat -= nb_treated

    # If there are leftovers (there should not be any)
    if list(iterable[treated:]):
        logging.error(
            'There are leftovers in `uniform_slice` this should not '
            'happen (iterable len: {}, max_slice_size: {})'.format(
                len(iterable), max_slice_size))
        yield iterable[treated:]


"""
    Helper function for pd.DataFrame
"""


def append_name(df, series, name):
    """Append a row at the end of a pd.DataFrame with the specified index `name`

    [description]
    :param df: DataFrame to append `series` to
    :type df: pd.DataFrame
    :param series: List of values to append to `df`,
                   df's column length and series's len must match
    :type series: pd.Series / list
    :param name: Index of the appent series
    :type name: str
    :returns: [description]
    :rtype: {pd.DataFrame}
    """
    series.name = name
    return df.append(series)


def ensure_indexes(df, indexes):
    """Ensure indexes exists in a DataFrame

    Reorders the indexes in df. The firsts wil be `indexes` and the
    rest in the same order as before
    :param df: dataframe to add indexes to
    :type df: pd.DataFrame
    :param indexes: List of indexes
    :type indexes: list
    :returns: DataFrame with at least `indexes` indexes
    :rtype: {pd.DataFrame}
    """
    tmp = df.loc[[i for i in df.index if i not in indexes]]
    return df.reindex(index=indexes).append(tmp)


def ensure_columns(df, columns):
    """Ensure columns exists in a DataFrame

    Reorders the columns in df. The firsts wil be `columns` and the
    rest in the same order as before
    :param df: dataframe to add columns to
    :type df: pd.DataFrame
    :param columns: List of columns
    :type columns: list
    :returns: DataFrame with at least `columns` columns
    :rtype: {pd.DataFrame}
    """
    tmp = [i for i in df.columns if i not in columns]
    return df.reindex(columns=(columns + tmp))


def set_name(df, name):
    """Add name to DataFrame in various places to ensure it appears
    when using pd2pptx

    [description]
    :param df: DataFrame to add name to
    :type df: pd.DataFrame
    :param name: name to set
    :type name: str
    :returns: dataframe with name set
    :rtype: {pd.DataFrame}
    """

    # TODO comment this

    if df.axes[0].name:
        df.axes[0].name = name
    else:
        df.axes[0].names = [name] + [''] * (len(df.axes[1].names) - 1)

    if df.axes[1].name:
        df.axes[1].name = ''
    else:
        df.axes[1].names = [''] * len(df.axes[1].names)

    return df


"""
    PowerPoint Generation
"""


def add_chart_slide(slide, df, chart_type, title_size=None, text_size=None,
                    ind=None, cols=None, data_labels=None, name=None,
                    legend_pos=None, font=None, datalabel_size=None,
                    hide_zeros=None):
    """
        Add a chart to a slide using a DataFrame
        TODO : create a function to create the ChartData from DF
        TODO : create a function to insert chart in placeholder or frame
        TODO : create a function to apply style to created chart
    """

    if datalabel_size is None:
        datalabel_size = text_size
    if chart_type == CHART.XY_SCATTER:
        data_labels = False

    # Select dataframe slice
    if ind is None:
        ind = (None, None)

    if cols is None:
        cols = (None, None)

    df = df.iloc[ind[0]:ind[1], cols[0]:cols[1]]

    if chart_type == CHART.XY_SCATTER:
        # Create ChartData
        chart_data = XyChartData()
        df = df.fillna(0)
        for ind, row in df.iterrows():
            s = chart_data.add_series(ind)
            s.add_data_point(row[0], row[1])
    else:
        # Create ChartData
        chart_data = ChartData()
        chart_data.categories = list(df.columns)

        # In order to not print 0 data_labels
        if hide_zeros:
            df = df.replace(0, np.nan)
        df = df.fillna('')
        for ind, row in df.iterrows():
            chart_data.add_series(ind, row)

    # Insert chart in slide or placeholder
    try:
        shp = slide.insert_chart(chart_type, chart_data)
    except AttributeError:
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        shp = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data
        )

    # Define chart style

    if name:
        title = shp.chart.chart_title.text_frame
        title.text = name
        for p in title.paragraphs:
            if font:
                p.font.name = font
            if title_size:
                p.font.size = Pt(title_size)

    # Legend
    shp.chart.has_legend = True
    shp.chart.legend.include_in_layout = False
    if legend_pos:
        shp.chart.legend.position = legend_pos

    # Text sizes
    if text_size is not None:
        text_size = Pt(text_size)
        # Lower indexes
        try:
            shp.chart.value_axis.tick_labels.font.size = text_size
        except ValueError:
            ()
        for plot in shp.chart.plots:
            # Legend
            plot.chart.legend.font.size = text_size
            # Indexes
            if chart_type not in [CHART.DOUGHNUT, CHART.PIE]:
                plot.chart.category_axis.tick_labels.font.size = text_size

            try:
                plot.has_data_labels = data_labels
                if plot.has_data_labels:
                    # Numbers in the plot (on the bars)
                    plot.data_labels.font.size = Pt(datalabel_size)
            except AttributeError:
                logging.info('No datalabels available for {} chart'.format(chart_type._member_name))
    return shp


def map_cells(func, table, index=None):
    """Apply a function on each cell in a pptx.Table

    Mostly for styling
    :param func: function to apply
    :type func: cell -> row_id -> col_id -> void
    :param table: Table to iter
    :type table: pptx.Table
    :param index: Use integer as index instead of axis name, defaults to False
    :type index: bool, optional
    :returns: the input table
    :rtype: {pptx.Table}
    """
    if index is None:
        index = False
    for i, j in product(range(len(table.rows)), range(len(table.columns))):
        if index:
            func(table.cell(i, j), j, i)
        else:
            func(table.cell(i, j),
                 table.cell(0, j).text_frame.text,
                 table.cell(i, 0).text_frame.text)
    return table


def color_rows_except_1st_col(rows, color=None):
    """Functor Color row with color

    To use with `map_cells`
    :param rows: list of rows to color
    :type rows: list
    :param color: color, defaults to None
    :type color: pptx.enum.dml.MSO_THEME_COLOR, optional
    :returns: function to set pptx.table._Cell style
    :rtype: {funtion}
    """
    if color is None:
        color = COLOR.ACCENT_3

    def tmp(cell, col, row):
        if col != 0 and row in rows:
            cell.fill.solid()
            cell.fill.fore_color.theme_color = color
    return tmp


def color_cell_diff(values, color=None):
    """Functor Color row with color

    To use with `map_cells`
    :param values: values
    :type values: list
    :param color: color, defaults to None
    :type color: pptx.enum.dml.MSO_THEME_COLOR, optional
    :returns: function to set pptx.table._Cell style
    :rtype: {funtion}
    """
    if color is None:
        color = COLOR.ACCENT_4

    def tmp(cell, col, row):
        try:
            value = int(cell.text_frame.text)
        except ValueError:
            value = 0
        if row == 'Total' or row == '' or col == '':
            return
        if values[row] != value:
            cell.fill.solid()
            cell.fill.fore_color.theme_color = color
            cell.fill.fore_color.brightness = 0.5
    return tmp


class PptxWriter(object):
    """Helper class to use pptx-python package"""

    # Format for layout name:
    # (\d)_(\w(_\w)*)(_c|r)?:
    # \1 : number of \2 available
    # \2 : object type (chart, table)
    # \3 : row or column
    layouts = {'title': 0, '2_cha': 1, '2_tab_cha_r': 2,
               '4_cha': 3, '1_tab_cha_c': 4, '1_tab_cha_c_center': 5,
               '1_tab_cha_r': 6, '1_tab': 7, '1_cha': 8,
               '2_tab_r': 9, '2_cha_c': 10, '3_1_cha': 11,
               'ending': [12, 13], '3_cha': 14, '2_1_cha': 15}
    default_layout_index = layouts['2_tab_cha_r']

    def __init__(self, prs):
        super(PptxWriter, self).__init__()
        self.prs = prs
        # Actual slide infos
        self.slide = None
        self.title = None
        self.subtitle = None
        self.page_number = 1
        # Number of written elements
        self.placeholders = None
        self.nb_written = 0
        # Used for content slides
        self.layout_index = None
        self.set_layout_index(PptxWriter.default_layout_index)

    def _add_slide(self, layout_index):
        """Add a basic new slide using layout_index

        [description]
        :param layout_index: [description]
        :type layout_index: [type]
        :returns: [description]
        :rtype: {[type]}
        """
        # remove every empty placeholder
        if self.slide:
            for i in self.slide.placeholders:
                # Empty placeholders have empty text_frame
                if i.has_text_frame:
                    if not i.text:
                        i = i._element
                        i.getparent().remove(i)

        self.slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[layout_index])
        try:
            self.slide.placeholders[2].text = self.page_number
        except KeyError:
            ()
        self.page_number += 1
        return self.slide

    def section(self, name, write=True):
        """Add a new section slide

        :param name: name of the section
        :type name: str
        :param write: write this slide, defaults to True
        :type write: bool, optional
        """
        self.title = name
        self.subtitle = None
        self.set_slide_name(None)
        if write:
            self._add_slide(PptxWriter.layouts['title'])
            self.slide.shapes.title.text = self.title

    def subsection(self, name, write=True):
        """Add a new subsection slide

        [description]
        :param name: name of the subsection
        :type name: str
        :param write: write this slide, defaults to True
        :type write: bool, optional
        """
        self.subtitle = name
        self.set_slide_name(None)
        self._add_slide(PptxWriter.layouts['title'])
        if write:
            self.slide.shapes.title.text = self.title
            self.slide.placeholders[14].text = self.subtitle

    def set_layout_index(self, layout_index):
        """Change layout_index for content_slide

        Updates `self.placeholders`
        :param layout_index: index of the layout to be used
        :type layout_index: int
        """
        if layout_index == self.layout_index:
            return
        self.layout_index = layout_index
        placeholders = self.prs.slide_layouts[self.layout_index].placeholders
        self.placeholders = [shp.placeholder_format.idx for shp
                             in placeholders]
        self.placeholders = sorted(self.placeholders)
        # Removing the first 2 placeholders which are the title and subtitle
        self.placeholders = self.placeholders[2:]

    def content_slide(self, name=None, layout_index=None):
        # TODO Rename to say that it add a slide
        """Add a new content slide

        [description]
        :param name: name of the slide, defaults to None
        :type name: str, optional
        :param layout_index: index of the slide type, defaults to None
        :type layout_index: PptxWriter.layouts, optional
        """
        if layout_index is not None:
            self.set_layout_index(layout_index)
        self.list_placeholders(self.layout_index)
        if name:
            self.set_slide_name(name)

        self._add_slide(self.layout_index)
        self.slide.shapes.title.text = self.slide_name
        self.slide.placeholders[13].text += ' - '.join(
            filter(None, [self.title, self.subtitle]))
        self.nb_written = 0

    def next_placeholder(self):
        """Return the next placeholder index

        [description]
        :returns: placeholder index
        :rtype: {int}
        """
        # If there are no more palceholders go to the next slide
        if (self.nb_written >= len(self.placeholders)) \
           or (self.slide is None):
            self.content_slide()

        res = self.slide.placeholders[self.placeholders[self.nb_written]]
        self.nb_written += 1
        return res

    def set_slide_name(self, name):
        self.slide_name = name

    def add_table_and_chart(
            self, df, df_chart=None, index=None, chart_type=None, ind=None, cols=None,
            name=None, hide_zeros=None, data_labels=None, legend_pos=None,
            font=None, text_align=None, text_size=None, title_size=None,
            datalabel_size=None, height=None):
        """Add a table and a chart to the presentation

        [description]
        :param df: [description]
        :type df: [type]
        :returns: tuple containing written table and chart
        :rtype: {(pptx.shapes.graphfrm.GraphicFrame,
                  pptx.shapes.graphfrm.GraphicFrame)}
        """

        if df_chart is None:
            df_chart = df

        table = self.add_table(
            df, text_size=text_size, index=index, text_align=text_align,
            hide_zeros=hide_zeros)

        chart = self.add_chart(
            df_chart, chart_type=chart_type, text_size=text_size, ind=ind,
            cols=cols, name=name, legend_pos=legend_pos,
            font=font, datalabel_size=datalabel_size,
            data_labels=data_labels, hide_zeros=hide_zeros)

        return table, chart

    def add_table(self, df, index=None, hide_zeros=None, text_size=None,
                  text_align=None, row_height=None, **kwargs):
        """Add a table in the presentation

        [description]
        :param df: data to use
        :type df: pd.DataFrame
        :param text_size: size of text content, defaults to None
        :type text_size: float, optional
        :param index: display the indexes, defaults to True
        :type index: boolean, optional
        :param row_height: row height, defaults to None
        :type row_height: float, optional
        :param text_align: text alignement, defaults to None
        :type text_align: pptx.enum.text.PP_ALIGN optional
        :param hide_zeros: do not print 0, defaults to True
        :type hide_zeros: boolean, optional
        :param **kwargs: UNUSED HACK to collect extra arguments
        :type **kwargs: dict
        :returns: written table
        :rtype: {pptx.shapes.graphfrm.GraphicFrame}
        """

        if text_align is None:
            text_align = PP_ALIGN.RIGHT
        if hide_zeros is None:
            hide_zeros = True
        if index is None:
            index = True

        df = df.apply(partial(round, ndigits=2))

        table = pd2pptx.df_to_table(
            self.next_placeholder(), df, text_size=text_size, index=index,
            text_align=text_align, row_height=row_height,
            hide_zeros=hide_zeros)

        pd2pptx.set_table_style(table, TableStyle.LightStyle2Accent1)

        # Add color to first col
        def tmp(cell, col, row):
            if col == 0 and row != 0:
                cell.fill.solid()
                cell.fill.fore_color.theme_color = COLOR.ACCENT_2
                cell.fill.fore_color.brightness = 0.3
            if row == 0:
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(255, 255, 255)
        map_cells(tmp, table.table, index=True)

        return table

    def add_chart(
            self, df, chart_type=None, ind=None, cols=None,
            name=None, hide_zeros=None, data_labels=None,
            legend_pos=None, font=None, text_size=None,
            title_size=None, datalabel_size=None,
            show_percent=None, chart_style=None,
            minor_unit=None, scale=None):
        """Add a chart in the Presentation

        [description]
        :param df: data to use
        :type df: pd.DataFrame
        :param chart_type: Type of chart, defaults to BAR_STACKED_100
        :type chart_type: pptx.enum.chart.XL_CHART_TYPE, optional
        :param hide_zeros: do not print 0, defaults to True
        :type hide_zeros: boolean, optional
        :param text_size: size of text content, defaults to 10.5
        :type text_size: float, optional
        :param ind: dataframe slice to use, defaults to None
        :type ind: (int or None, int or None), optional
        :param cols: dataframe slice to use, defaults to None
        :type cols: (int or None, int or None), optional
        :param name: name of the chart, defaults to None
        :type name: str, optional
        :param legend_pos: legend position, defaults to None
        :type legend_pos: pptx.enum.chart.XL_LEGEND_POSITION, optional
        :param font: font name to use, defaults to None
        :type font: str, optional
        :param title_size: title size, defaults to None
        :type title_size: float, optional
        :param datalabel_size: datalabel size, defaults to None
        :type datalabel_size: float, optional
        :param data_labels: print the datalabels, defaults to None
        :type data_labels: boolean, optional
        :returns: written chart
        :rtype: {pptx.shapes.graphfrm.GraphicFrame}
        """

        if chart_type is None:
            chart_type = CHART.BAR_STACKED_100
        if text_size is None:
            text_size = 10.5
        if hide_zeros is None:
            hide_zeros = True
        if show_percent is None:
            show_percent = True
        if chart_style is None:
            chart_style = 2
        if minor_unit is None:
            minor_unit = 1
        if scale is None:
            scale = (None, None)
        if data_labels is None:
            data_labels = True

        use_datalabels = (data_labels or (show_percent and chart_type == CHART.DOUGHNUT))

        df = df.apply(partial(round, ndigits=2))
        chart = add_chart_slide(
            self.next_placeholder(), df, chart_type, ind=ind, cols=cols,
            name=name, hide_zeros=hide_zeros, data_labels=use_datalabels,
            legend_pos=legend_pos, font=font, text_size=text_size,
            title_size=title_size, datalabel_size=datalabel_size)

        if use_datalabels and (chart and chart_type not in [CHART.XY_SCATTER]):
            showPercent_elt = chart.chart._chartSpace.xpath(
                '//c:showVal')[0]
            showPercent_elt.attrib['val'] = str(int(data_labels))

        if chart_type in [CHART.DOUGHNUT]:
            # Set holeSize for Doughnut Charts
            holeSize_elt = chart.chart._chartSpace.xpath(
                'c:chart/c:plotArea/c:doughnutChart/c:holeSize')[0]
            holeSize_elt.attrib['val'] = '70'

            firstSliceAng_elt = chart.chart._chartSpace.xpath(
                'c:chart/c:plotArea/c:doughnutChart/c:firstSliceAng')[0]
            firstSliceAng_elt.attrib['val'] = '180'

            # Set DataLabels

            if use_datalabels and show_percent:
                try:
                    showPercent_elt = chart.chart._chartSpace.xpath(
                        'c:chart/c:plotArea/c:doughnutChart/c:dLbls/c:showPercent')[0]
                    showPercent_elt.attrib['val'] = '1'
                except IndexError:
                    showPercent_elt = chart.chart._chartSpace.xpath(
                        'c:chart/c:plotArea/c:doughnutChart/c:dLbls')[0]
                    print(dir(showPercent_elt))
                    showPercent_elt.attrib['val'] = '1'

        chart.chart.chart_style = chart_style
        try:
            chart.chart.value_axis.minor_unit = minor_unit
            if scale[0]:
                chart.chart.value_axis.minimum_scale = scale[0]
            if scale[1]:
                chart.chart.value_axis.maximum_scale = scale[1]
        except ValueError:
            ()  # No value axis

        return chart

    def add_ending_slides(self):
        for i in PptxWriter.layouts['ending']:
            self._add_slide(i)

    def save(self, output_file):
        """Save the Presentation to file `output_file`

        If `output_file` already exists overwrite it,
        ensures that path exist
        :param output_file: file to write to
        :type output_file: str
        """
        os.makedirs(os.path.dirname(output_file), exist_ok=True)
        self.prs.save(output_file)

    def list_placeholders(self, layout_index=None):
        """Debugging tool to list placeholders in a template layout

        :param layout_index: index of the layout, defaults to None
        :type layout_index: int, optional
        """
        if layout_index is None:
            layout_index = self.layout_index
        logging.debug('List placeholders in {}'.format(layout_index))
        for shape in self.prs.slide_layouts[layout_index].placeholders:
            logging.debug('{} {}'.format(
                shape.placeholder_format.idx, shape.name))


def add_table_and_chart_per_pack(
        writer, name, df, row, args, max_col,
        fn_map_cell=None, fn_map_cell_index=None,
        last_row=None):
    """Function to split large tables to be written on a pptx.Presentation

    [description]
    :param writer: presentation writer
    :type writer: PptxWriter
    :param name: name of slide
    :type name: str
    :param df: data to use
    :type df: pd.DataFrame
    :param row: number of row in the dataframe
    :type row: int
    :param args: arguments for table and chart styling
    :type args: dict
    :param max_col: max number of columns per written table
    :type max_col: int
    :param fn_map_cell: function to apply to table, defaults to None
    :type fn_map_cell: fun, optional
    :param fn_map_cell_index: , defaults to None
    :type fn_map_cell_index: [type], optional
    """
    # If there are too many rows then use a full slide per table slice
    columns_pack = list(uniform_slice(df.columns, max_col))

    if row > 7:
        layout = PptxWriter.layouts['1_tab']
    else:
        layout = PptxWriter.layouts['2_tab_r']

    if len(columns_pack) == 1:
        layout = PptxWriter.layouts['2_tab_r']

    writer.content_slide(name, layout)

    tables = []
    for cols in columns_pack:
        table = writer.add_table(df[cols], **args)
        tables.append(table)
        if fn_map_cell:
            map_cells(fn_map_cell, table.table, index=fn_map_cell_index)
        if last_row is not None:
            table.table.last_row = last_row

    # TODO : if there is room for the chart to be with the last
    #        table use this space
    if len(columns_pack) > 1:
        writer.content_slide(name, PptxWriter.layouts['1_cha'])
    else:
        if 'chart_type' in args and args['chart_type'] == CHART.RADAR:
            args['legend_pos'] = LEGEND_POS.RIGHT
    chart = writer.add_chart(df, **args)
    return tables, chart


def add_table_per_pack(
        writer, name, df, row, args, max_col,
        fn_map_cell=None, fn_map_cell_index=None,
        last_row=None):
    """Function to split large tables to be written on a pptx.Presentation

    [description]
    :param writer: presentation writer
    :type writer: PptxWriter
    :param name: name of slide
    :type name: str
    :param df: data to use
    :type df: pd.DataFrame
    :param row: number of row in the dataframe
    :type row: int
    :param args: arguments for table and chart styling
    :type args: dict
    :param max_col: max number of columns per written table
    :type max_col: int
    :param fn_map_cell: function to apply to table, defaults to None
    :type fn_map_cell: fun, optional
    :param fn_map_cell_index: , defaults to None
    :type fn_map_cell_index: [type], optional
    """
    # If there are too many rows then use a full slide per table slice
    columns_pack = list(uniform_slice(df.columns, max_col))

    if row > 7:
        layout = PptxWriter.layouts['1_tab']
    else:
        layout = PptxWriter.layouts['2_tab_r']

    if len(columns_pack) == 1:
        layout = PptxWriter.layouts['2_tab_r']

    writer.content_slide(name, layout)

    tables = []
    for cols in columns_pack:
        table = writer.add_table(df[cols], **args)
        tables.append(table)
        if fn_map_cell:
            map_cells(fn_map_cell, table.table, index=fn_map_cell_index)
        if last_row is not None:
            table.table.last_row = last_row

    return tables


def age_pyramid(
        df, writer, to_count_col, age_col, gender_col, genders=None,
        age_slice_width=None, title=None, extremums=None, open_slice=None, args=None):
    # TODO : Comment
    if age_slice_width is None:
        age_slice_width = 5
    if genders is None:
        genders = sorted(set(df[gender_col].dropna()))
    if title is None:
        title = 'Pyramide des âges'
    if open_slice is None:
        open_slice = True
    if extremums is None:
        extremums = (None, None)

    default_args = {
        'chart_type': CHART.BAR_STACKED, 'name': title,
        'title_size': 21.6,  'data_labels': True, 'datalabel_size': 12}
    if args:
        default_args.update(args)

    tmp = df.drop_duplicates(to_count_col).dropna(subset=[age_col])
    logging.info('Hard coded code')
    extremums = (24, 69)

    min_age, max_age = df[age_col].min(), df[age_col].max()
    if extremums:
        min_age = extremums[0] or min_age
        max_age = extremums[1] or max_age

    if min_age % 5 == 0:
        min_age -= 1

    lower_fifth = lambda x: x - (x % 5)
    min_age, max_age = lower_fifth(min_age) + 5, lower_fifth(max_age)

    age_slice = np.arange(min_age - age_slice_width, max_age + age_slice_width + 1, age_slice_width)


    ttmp = tmp.groupby([pd.cut(tmp[age_col], age_slice), gender_col])
    pyramid = ttmp[to_count_col].count().unstack()
    pyramid.index = pyramid.index.map(lambda x: str(x)[1:-1])

    if len(tmp) != pyramid.sum().sum():
        logging.error('Pyramid ill defined')
        logging.error('Persons: {}, in the pyramid: {}'.format(
            len(tmp), pyramid.sum().sum()))
        logging.error('Min age : {} Max age : {}'.format(df[age_col].min(), df[age_col].max()))
        logging.error(age_slice)


    slices = [', '.join(map(str, (map(int, age_slice[i - 1:i + 1])))) for i in range(1, len(age_slice))]
    pyramid = ensure_indexes(pyramid, slices)
    pyramid.sort_index(ascending=True, inplace=True)

    if open_slice and len(pyramid.index) > 1:
        to_update = {
            pyramid.index[-1]: pyramid.index[-1].split(', ')[0] + ' et plus',
            pyramid.index[0]: pyramid.index[0].split(', ')[1] + ' et moins'}

        pyramid.index = pyramid.index.map(
            lambda x: to_update[x] if x in to_update else x)
    # Minus a gender so both can appear on each side of the chart
    if len(pyramid) > 0:
        pyramid.iloc[:, 0] = -pyramid.iloc[:, 0]
        pyramid = ensure_columns(pyramid, genders)
        pyramid = pyramid.T

        # Styling for age pyramid
        chart = writer.add_chart(
            pyramid, **default_args)
        # The abscys is put to the left
        chart.chart.category_axis.tick_label_position = TICK_LABEL_POS.LOW
        # The bars are on the same level
        chart.chart.overlap = '100'
        # Do not print the minus on the labels
        for p in chart.chart.plots:
            p.data_labels.number_format = '# ##0; # ##0'
        for s in chart.chart.series:
            s.invert_if_negative = False
        return chart
    return None
